import os
from pathlib import Path

# PDF extraction
import pdfplumber
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# DOCX
from docx import Document
import docx2txt

# HTML
from bs4 import BeautifulSoup

# PPTX
from pptx import Presentation

# RTF
try:
    import striprtf
except ImportError:
    striprtf = None

# Requests for URL fetching
import requests


def pdf_to_text_plumber(path_pdf: str) -> str:
    text = []
    with pdfplumber.open(path_pdf) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return "\n".join(text)


def pdf_to_text_pymupdf(path_pdf: str) -> str:
    if fitz is None:
        raise ImportError("PyMuPDF is not installed")
    doc = fitz.open(path_pdf)
    text = []
    for page in doc:
        text.append(page.get_text("text"))
    return "\n".join(text)


def docx_to_text_python_docx(path_docx: str) -> str:
    doc = Document(path_docx)
    return "\n".join([p.text for p in doc.paragraphs])


def docx_to_text2(path_docx: str) -> str:
    return docx2txt.process(path_docx)


def html_to_text_string(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    lines = [line.strip() for line in soup.get_text().splitlines()]
    return "\n".join([l for l in lines if l])


def html_to_text(path_html: str) -> str:
    with open(path_html, "r", encoding="utf-8") as f:
        html = f.read()
    return html_to_text_string(html)


def url_to_text(url: str) -> str:
    resp = requests.get(url, timeout=10)
    resp.raise_for_status()
    return html_to_text_string(resp.text)


def pptx_to_text(path_pptx: str) -> str:
    prs = Presentation(path_pptx)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def rtf_to_text(path_rtf: str) -> str:
    if striprtf is None:
        raise ImportError("striprtf is not installed")
    with open(path_rtf, "r", encoding="utf-8") as f:
        rtf = f.read()
    return striprtf.rtf_to_text(rtf)


def txt_to_text(path_txt: str) -> str:
    with open(path_txt, "r", encoding="utf-8") as f:
        return f.read()


def ingest_to_text(path: str) -> str:
    """
    Detect file type by extension and return extracted plain text.
    Supported: .pdf, .docx, .html/.htm, .txt, .pptx, .rtf, http/https URLs
    """
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if path.startswith(('http://', 'https://')):
            return url_to_text(path)
        if ext == '.pdf':
            try:
                return pdf_to_text_plumber(path)
            except Exception:
                return pdf_to_text_pymupdf(path)
        if ext == '.docx':
            try:
                return docx_to_text_python_docx(path)
            except Exception:
                return docx_to_text2(path)
        if ext in ('.html', '.htm'):
            return html_to_text(path)
        if ext == '.txt':
            return txt_to_text(path)
        if ext == '.pptx':
            return pptx_to_text(path)
        if ext == '.rtf':
            return rtf_to_text(path)
        raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        raise RuntimeError(f"Error processing {path}: {e}")


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Extract plain text from various document formats.')
    parser.add_argument('input', help='Path or URL of the document')
    parser.add_argument('-o', '--output', help='Output text file (optional)')
    args = parser.parse_args()
    txt = ingest_to_text(args.input)
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as fout:
            fout.write(txt)
        print(f"Text written to {args.output}")
    else:
        print(txt)
