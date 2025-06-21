import os
import argparse
from pathlib import Path
import mimetypes
import logging
import re
import json
import csv

# Logging configuration
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# PDF extraction - Using both libraries as fallback
import pdfplumber
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
    logger.warning("PyMuPDF not installed. PDF extraction may be less robust")

# DOCX - Using both methods
import docx
try:
    import docx2txt
except ImportError:
    docx2txt = None
    logger.warning("docx2txt not installed. Using only python-docx for DOCX")

# PowerPoint files
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logger.warning("python-pptx not installed. PPTX files not supported")

# RTF files
try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None
    logger.warning("striprtf not installed. RTF files not supported")


# OCR capabilities (optional)
try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("OCR not available. Image files not supported")

# Email processing
try:
    import email
    from email import policy
    from email.parser import BytesParser
except ImportError:
    logger.warning("Email module not found. EML files not supported")

# Excel processing
try:
    import pandas as pd
except ImportError:
    pd = None
    logger.warning("pandas not installed. XLS/XLSX files not supported")


def clean_text(text: str) -> str:
    """
    Clean extracted text by removing control characters and normalizing whitespace.

    Parameters:
    ----------
    text : str
        The raw text extracted from the document.

    Returns:
    -------
    str
        Cleaned text with control characters removed, whitespace normalized, and leading/trailing spaces stripped.
    """
    if not text:
        return ""
    
    # Remove control characters (except tabs and newlines)
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Strip leading/trailing spaces
    return text.strip()


def extract_pdf_text(path: str) -> str:
    """
    Extract text from PDF using multiple methods with fallback.

    Parameters:
    ----------
    path : str
        Path to the PDF file.

    Returns:
    -------
    str
        Extracted text from the PDF file.
    """
    text = ""
    
    # Method 1: pdfplumber (better for tables and formatting)
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return clean_text(text)
    except Exception as e:
        logger.warning(f"pdfplumber failed: {str(e)}")
    
    # Method 2: PyMuPDF (fallback)
    if fitz:
        try:
            doc = fitz.open(path)
            text = []
            for page in doc:
                text.append(page.get_text("text", sort=True))
            return clean_text("\n".join(text))
        except Exception as e:
            logger.warning(f"PyMuPDF failed: {str(e)}")
    
    raise RuntimeError("All PDF extraction methods failed")


def extract_docx_text(path: str) -> str:
    """
    Extract text from DOCX using both python-docx and docx2txt as fallback.

    Parameters:
    ----------
    path : str
        Path to the DOCX file.

    Returns:
    -------
    str
        Extracted text from the DOCX file.
    """
    text = ""
    
    # Method 1: python-docx (better for structured text)
    try:
        doc = docx.Document(path)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\t"
                text += "\n"
    except Exception as e:
        logger.warning(f"python-docx partially failed: {str(e)}")
    
    # Method 2: docx2txt (better for raw extraction)
    if docx2txt:
        try:
            alt_text = docx2txt.process(path)
            # Combine both results avoiding duplicates
            if alt_text and alt_text not in text:
                text += "\n" + alt_text
        except Exception as e:
            logger.warning(f"docx2txt failed: {str(e)}")
    
    return clean_text(text) if text else ""


def extract_pptx_text(path: str) -> str:
    """
    Extract text from PowerPoint PPTX files.

    Parameters:
    ----------
    path : str
        Path to the PPTX file.

    Returns:
    -------
    str
        Extracted text from the PPTX file.
    """
    if Presentation is None:
        raise RuntimeError("PPTX processing not available (python-pptx not installed)")
    
    try:
        prs = Presentation(path)
        text = []
        
        # Process each slide's shapes and tables
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text_frame.text.strip())
                        text.append("\t".join(row_text))
                
                # Handle grouped shapes
                if shape.shape_type == 6:  # Group shape type
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text"):
                            text.append(subshape.text.strip())
        
        # Extract notes text from slides if any
        for slide in prs.slides:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                text.append(notes_slide.notes_text_frame.text.strip())
        
        return clean_text("\n".join(filter(None, text)))
    
    except Exception as e:
        raise RuntimeError(f"Failed to extract PPTX text: {str(e)}")


def extract_rtf_text(path: str) -> str:
    """
    Extract text from RTF files.

    Parameters:
    ----------
    path : str
        Path to the RTF file.

    Returns:
    -------
    str
        Extracted text from the RTF file.
    """
    if rtf_to_text is None:
        raise RuntimeError("RTF processing not available (striprtf not installed)")

    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            return clean_text(text)
    except Exception as e:
        raise RuntimeError(f"Failed to extract RTF text: {str(e)}")
    

def extract_image_text(path: str) -> str:
    """
    Extract text from images using OCR.

    Parameters:
    ----------
    path : str
        Path to the image file.

    Returns:
    -------
    str
        Extracted text from the image file using OCR.
    """
    if not OCR_AVAILABLE:
        raise RuntimeError("OCR not available. Install pytesseract and PIL")
    
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        return clean_text(text)
    except Exception as e:
        raise RuntimeError(f"OCR failed: {str(e)}")
    

def extract_eml_text(path: str) -> str:
    """
    Extract text from .eml email files.

    Parameters:
    ----------
    path : str
        Path to the .eml file.

    Returns:
    -------
    str
        Extracted text from the .eml file.
    """
    try:
        with open(path, 'rb') as f:
            msg = BytesParser(policy=policy.default).parse(f)
        text = ""
        
        # Extract text/plain parts from multipart emails
        if msg.is_multipart():
            for part in msg.iter_parts():
                if part.get_content_type() == 'text/plain':
                    text += part.get_content()
        else:
            if msg.get_content_type() == 'text/plain':
                text = msg.get_content()
        
        return clean_text(text)
    except Exception as e:
        raise RuntimeError(f"Error processing email: {str(e)}")
    

def extract_csv_text(path: str) -> str:
    """
    Extract text from CSV files.

    Parameters:
    ----------
    path : str
        Path to the CSV file.

    Returns:
    -------
    str
        Extracted text from the CSV file as rows joined by pipe (|).
    """
    try:
        rows = []
        with open(path, 'r', encoding='utf-8', errors='replace') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                rows.append(" | ".join(row))
        return clean_text("\n".join(rows))
    except Exception as e:
        raise RuntimeError(f"Error processing CSV: {str(e)}")


def extract_json_text(path: str) -> str:
    """
    Extract text from JSON files with pretty formatting.

    Parameters:
    ----------
    path : str
        Path to the JSON file.

    Returns:
    -------
    str
        Extracted and pretty-printed JSON text.
    """
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            data = json.load(f)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        return clean_text(text)
    except Exception as e:
        raise RuntimeError(f"Error processing JSON: {str(e)}")


def extract_excel_text(path: str) -> str:
    """
    Extract text from Excel files (.xls, .xlsx) by reading all sheets.

    Parameters:
    ----------
    path : str
        Path to the Excel file.

    Returns:
    -------
    str
        Extracted text from all sheets concatenated.
    """
    if pd is None:
        raise RuntimeError("pandas not installed. Cannot process Excel files.")
    
    try:
        # Read all sheets as a dictionary of DataFrames
        df_dict = pd.read_excel(path, sheet_name=None)
        text_parts = []
        for sheet_name, df in df_dict.items():
            text_parts.append(f"--- Sheet: {sheet_name} ---\n")
            # Convert DataFrame to string without the index column
            text_parts.append(df.to_string(index=False))
            text_parts.append("\n")
        return clean_text("\n".join(text_parts))
    except Exception as e:
        raise RuntimeError(f"Error processing Excel: {str(e)}")


def file_to_text(path: str) -> str:
    """
    Detect file type by MIME or extension and extract text using appropriate method.

    Supports: PDF, DOCX, TXT, PPTX, RTF, images, emails, CSV, JSON, Excel

    Parameters:
    ----------
    path : str
        Path to the file.

    Returns:
    -------
    str
        Extracted text content.
    """
    mime_type, _ = mimetypes.guess_type(path)

    # Use MIME type to decide extraction method
    if mime_type:
        if mime_type == 'application/pdf':
            return extract_pdf_text(path)
        elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                           'application/msword']:
            return extract_docx_text(path)
        elif mime_type == 'text/plain':
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                return clean_text(f.read())
        elif mime_type == 'application/vnd.ms-powerpoint':
            return extract_pptx_text(path)
        elif mime_type == 'application/rtf':
            return extract_rtf_text(path)
        elif mime_type.startswith('image/'):
            return extract_image_text(path)
        elif mime_type == 'message/rfc822':
            return extract_eml_text(path)
        elif mime_type == 'text/csv':
            return extract_csv_text(path)
        elif mime_type == 'application/json':
            return extract_json_text(path)
        elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                           'application/vnd.ms-excel']:
            return extract_excel_text(path)

    # Fallback on file extension
    ext = Path(path).suffix.lower()
    if ext == '.pdf':
        return extract_pdf_text(path)
    elif ext in ['.docx', '.doc']:
        return extract_docx_text(path)
    elif ext == '.txt':
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return clean_text(f.read())
    elif ext == '.pptx':
        return extract_pptx_text(path)
    elif ext == '.ppt':
        raise RuntimeError("Old .ppt format not supported natively")
    elif ext == '.rtf':
        return extract_rtf_text(path)
    elif ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']:
        return extract_image_text(path)
    elif ext == '.eml':
        return extract_eml_text(path)
    elif ext == '.csv':
        return extract_csv_text(path)
    elif ext == '.json':
        return extract_json_text(path)
    elif ext in ['.xls', '.xlsx']:
        return extract_excel_text(path)

    raise ValueError(f"Unsupported file type: {path}")


if __name__ == '__main__':

    parser = argparse.ArgumentParser(description='Text extraction from various document formats')
    parser.add_argument('input', help='Path of the document')
    parser.add_argument('-o', '--output', help='Output text file (optional)')
    parser.add_argument('-v', '--verbose', action='store_true', help='Verbose mode')

    args = parser.parse_args()

    if args.verbose:
        logger.setLevel(logging.DEBUG)

    try:
        txt = file_to_text(args.input)
        
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as fout:
                fout.write(txt)
            print(f"Text written to {args.output}")
        else:
            print(txt)
    except Exception as e:
        logger.error(f"Error processing file: {str(e)}")
        exit(1)

