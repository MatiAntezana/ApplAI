import os
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass

# PDF extraction
import pdfplumber
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# DOCX
from docx import Document
import docx2txt

# HTML
from bs4 import BeautifulSoup

# PPTX
from pptx import Presentation

# RTF
try:
    import striprtf
except ImportError:
    striprtf = None

# Requests for URL fetching
import requests

# Text processing
import unicodedata
from collections import Counter

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@dataclass
class TextMetrics:
    """Metrics for evaluating text quality"""
    char_count: int
    word_count: int
    line_count: int
    avg_line_length: float
    whitespace_ratio: float
    special_char_ratio: float
    repeated_char_ratio: float


class CVTextProcessor:
    """Enhanced text processor specifically designed for CV documents"""
    
    def __init__(self):
        # Common CV section headers in multiple languages
        self.cv_sections = {
            'personal_info': [
                'personal information', 'contact', 'contact information', 'personal details',
                'información personal', 'contacto', 'datos personales'
            ],
            'objective': [
                'objective', 'career objective', 'professional summary', 'summary',
                'objetivo', 'resumen profesional', 'perfil profesional'
            ],
            'experience': [
                'experience', 'work experience', 'professional experience', 'employment',
                'experiencia', 'experiencia laboral', 'historial laboral'
            ],
            'education': [
                'education', 'academic background', 'qualifications',
                'educación', 'formación académica', 'estudios'
            ],
            'skills': [
                'skills', 'technical skills', 'competencies', 'abilities',
                'habilidades', 'competencias', 'destrezas'
            ],
            'certifications': [
                'certifications', 'certificates', 'licenses',
                'certificaciones', 'certificados', 'licencias'
            ]
        }
        
        # Patterns for common formatting issues
        self.cleanup_patterns = [
            (r'\s+', ' '),  # Multiple spaces to single space
            (r'\n\s*\n\s*\n+', '\n\n'),  # Multiple newlines to double newline
            (r'[^\S\n]+$', '', re.MULTILINE),  # Trailing whitespace
            (r'^[^\S\n]+', '', re.MULTILINE),  # Leading whitespace
            (r'[•·▪▫◦‣⁃]\s*', '• '),  # Normalize bullet points
            (r'[\u2010-\u2015]', '-'),  # Normalize dashes
            (r'[\u2018\u2019]', "'"),  # Normalize single quotes
            (r'[\u201C\u201D]', '"'),  # Normalize double quotes
            (r'[\u2026]', '...'),  # Normalize ellipsis
        ]
        
        # Common junk patterns in CVs
        self.junk_patterns = [
            r'^\s*page\s+\d+\s*$',  # Page numbers
            r'^\s*\d+\s*/\s*\d+\s*$',  # Page x/y
            r'^[\s\-_=]{3,}$',  # Separator lines
            r'^\s*confidential\s*$',  # Confidential markers
            r'^\s*resume\s*$',  # Resume headers
            r'^\s*curriculum\s+vitae\s*$',  # CV headers
        ]

    def normalize_unicode(self, text: str) -> str:
        """Normalize Unicode characters for better consistency"""
        # Normalize to NFC form
        text = unicodedata.normalize('NFC', text)
        
        # Replace problematic Unicode characters
        replacements = {
            '\u00A0': ' ',  # Non-breaking space
            '\u2002': ' ',  # En space
            '\u2003': ' ',  # Em space
            '\u2009': ' ',  # Thin space
            '\u200B': '',   # Zero-width space
            '\u200C': '',   # Zero-width non-joiner
            '\u200D': '',   # Zero-width joiner
            '\uFEFF': '',   # Byte order mark
        }
        
        for old, new in replacements.items():
            text = text.replace(old, new)
        
        return text

    def clean_text(self, text: str) -> str:
        """Apply comprehensive text cleaning for better semantic consistency"""
        if not text:
            return ""
        
        # Normalize Unicode
        text = self.normalize_unicode(text)
        
        # Apply regex patterns
        for pattern, replacement, *flags in self.cleanup_patterns:
            flag = flags[0] if flags else 0
            text = re.sub(pattern, replacement, text, flags=flag)
        
        # Remove junk lines
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                cleaned_lines.append('')
                continue
                
            # Check if line matches junk patterns
            is_junk = any(re.match(pattern, line, re.IGNORECASE) for pattern in self.junk_patterns)
            
            if not is_junk:
                cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines)

    def preserve_structure(self, text: str) -> str:
        """Preserve important structural elements in CV text"""
        lines = text.split('\n')
        structured_lines = []
        
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                structured_lines.append('')
                continue
            
            # Detect section headers
            is_section_header = self.is_cv_section_header(line)
            
            # Add extra spacing before section headers (except first)
            if is_section_header and structured_lines and structured_lines[-1]:
                structured_lines.append('')
            
            structured_lines.append(line)
        
        return '\n'.join(structured_lines)

    def is_cv_section_header(self, line: str) -> bool:
        """Check if a line is likely a CV section header"""
        line_lower = line.lower().strip()
        
        # Check against known CV sections
        for section_type, headers in self.cv_sections.items():
            if any(header in line_lower for header in headers):
                return True
        
        # Check for common header patterns
        header_patterns = [
            r'^[A-Z\s]+$',  # All caps
            r'^[A-Z][a-z\s]*:?\s*$',  # Title case with optional colon
            r'^\w+\s*[-–—]\s*\w+',  # Word - Word format
        ]
        
        return any(re.match(pattern, line) for pattern in header_patterns)

    def calculate_metrics(self, text: str) -> TextMetrics:
        """Calculate text quality metrics"""
        if not text:
            return TextMetrics(0, 0, 0, 0.0, 0.0, 0.0, 0.0)
        
        lines = text.split('\n')
        words = text.split()
        
        char_count = len(text)
        word_count = len(words)
        line_count = len(lines)
        
        # Average line length
        non_empty_lines = [line for line in lines if line.strip()]
        avg_line_length = sum(len(line) for line in non_empty_lines) / len(non_empty_lines) if non_empty_lines else 0
        
        # Whitespace ratio
        whitespace_count = sum(1 for char in text if char.isspace())
        whitespace_ratio = whitespace_count / char_count if char_count > 0 else 0
        
        # Special character ratio
        special_char_count = sum(1 for char in text if not (char.isalnum() or char.isspace()))
        special_char_ratio = special_char_count / char_count if char_count > 0 else 0
        
        # Repeated character ratio (detect extraction artifacts)
        char_counts = Counter(text)
        repeated_chars = sum(count for char, count in char_counts.items() if count > len(text) * 0.05)
        repeated_char_ratio = repeated_chars / char_count if char_count > 0 else 0
        
        return TextMetrics(
            char_count=char_count,
            word_count=word_count,
            line_count=line_count,
            avg_line_length=avg_line_length,
            whitespace_ratio=whitespace_ratio,
            special_char_ratio=special_char_ratio,
            repeated_char_ratio=repeated_char_ratio
        )

    def validate_extraction_quality(self, text: str, file_path: str) -> Tuple[bool, List[str]]:
        """Validate the quality of extracted text"""
        issues = []
        metrics = self.calculate_metrics(text)
        
        # Check for common extraction issues
        if metrics.char_count < 100:
            issues.append("Text too short - possible extraction failure")
        
        if metrics.whitespace_ratio > 0.6:
            issues.append("Excessive whitespace - formatting issues detected")
        
        if metrics.special_char_ratio > 0.3:
            issues.append("High special character ratio - possible encoding issues")
        
        if metrics.repeated_char_ratio > 0.2:
            issues.append("High repeated character ratio - extraction artifacts detected")
        
        if metrics.avg_line_length < 5:
            issues.append("Lines too short - possible column extraction issues")
        
        # Check for garbled text patterns
        garbled_patterns = [
            r'[^\x00-\x7F]{10,}',  # Long sequences of non-ASCII
            r'(.)\1{5,}',  # Same character repeated 5+ times
            r'[A-Z]{20,}',  # Very long sequences of capitals
        ]
        
        for pattern in garbled_patterns:
            if re.search(pattern, text):
                issues.append(f"Garbled text detected: {pattern}")
        
        is_valid = len(issues) == 0
        return is_valid, issues


class EnhancedCVExtractor:
    """Enhanced CV text extractor with semantic consistency"""
    
    def __init__(self):
        self.processor = CVTextProcessor()
    
    def pdf_to_text_enhanced(self, path_pdf: str) -> str:
        """Enhanced PDF extraction with better text flow preservation"""
        text_chunks = []
        
        # Try pdfplumber first (better for structured documents)
        try:
            with pdfplumber.open(path_pdf) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract text with layout preservation
                    page_text = page.extract_text(
                        x_tolerance=3,
                        y_tolerance=3,
                        layout=True,
                        x_density=7.25,
                        y_density=13
                    )
                    
                    if page_text:
                        # Clean and process the page text
                        cleaned_text = self.processor.clean_text(page_text)
                        if cleaned_text.strip():
                            text_chunks.append(cleaned_text)
                            
        except Exception as e:
            logger.warning(f"pdfplumber failed for {path_pdf}: {e}")
            
            # Fallback to PyMuPDF
            if fitz is not None:
                try:
                    doc = fitz.open(path_pdf)
                    for page in doc:
                        # Extract text blocks to preserve layout
                        blocks = page.get_text("dict")["blocks"]
                        page_text = ""
                        
                        for block in blocks:
                            if "lines" in block:
                                for line in block["lines"]:
                                    line_text = ""
                                    for span in line["spans"]:
                                        line_text += span["text"]
                                    if line_text.strip():
                                        page_text += line_text + "\n"
                        
                        if page_text.strip():
                            cleaned_text = self.processor.clean_text(page_text)
                            text_chunks.append(cleaned_text)
                            
                except Exception as e2:
                    logger.error(f"Both PDF extractors failed for {path_pdf}: {e2}")
                    raise
        
        combined_text = "\n\n".join(text_chunks)
        return self.processor.preserve_structure(combined_text)

    def docx_to_text_enhanced(self, path_docx: str) -> str:
        """Enhanced DOCX extraction preserving document structure"""
        try:
            doc = Document(path_docx)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    # Preserve paragraph structure
                    text_parts.append(para_text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            combined_text = "\n".join(text_parts)
            cleaned_text = self.processor.clean_text(combined_text)
            return self.processor.preserve_structure(cleaned_text)
            
        except Exception as e:
            logger.warning(f"Enhanced DOCX extraction failed: {e}, falling back to simple extraction")
            return self.processor.clean_text(docx2txt.process(path_docx))

    def html_to_text_enhanced(self, html: str) -> str:
        """Enhanced HTML extraction preserving semantic structure"""
        soup = BeautifulSoup(html, "html.parser")
        
        # Remove unwanted elements
        for tag in soup(["script", "style", "meta", "link", "head"]):
            tag.decompose()
        
        # Extract text with preserved structure
        text_parts = []
        
        # Process common CV HTML structures
        cv_containers = soup.find_all(['div', 'section', 'article'], 
                                     class_=re.compile(r'(resume|cv|profile|experience|education|skills)', re.I))
        
        if cv_containers:
            for container in cv_containers:
                container_text = container.get_text().strip()
                if container_text:
                    text_parts.append(container_text)
        else:
            # Fallback to general text extraction
            text_parts.append(soup.get_text())
        
        combined_text = "\n\n".join(text_parts)
        cleaned_text = self.processor.clean_text(combined_text)
        return self.processor.preserve_structure(cleaned_text)

    def ingest_to_text_enhanced(self, path: str) -> Dict[str, any]:
        """
        Enhanced ingestion with quality metrics and validation
        Returns dict with text, metrics, and validation info
        """
        p = Path(path)
        ext = p.suffix.lower()
        
        try:
            # Extract text based on file type
            if path.startswith(('http://', 'https://')):
                resp = requests.get(path, timeout=10)
                resp.raise_for_status()
                raw_text = self.html_to_text_enhanced(resp.text)
            elif ext == '.pdf':
                raw_text = self.pdf_to_text_enhanced(path)
            elif ext == '.docx':
                raw_text = self.docx_to_text_enhanced(path)
            elif ext in ('.html', '.htm'):
                with open(path, "r", encoding="utf-8") as f:
                    html = f.read()
                raw_text = self.html_to_text_enhanced(html)
            elif ext == '.txt':
                with open(path, "r", encoding="utf-8") as f:
                    raw_text = self.processor.clean_text(f.read())
            elif ext == '.pptx':
                prs = Presentation(path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text.strip())
                raw_text = self.processor.clean_text("\n".join(text_parts))
            elif ext == '.rtf':
                if striprtf is None:
                    raise ImportError("striprtf is not installed")
                with open(path, "r", encoding="utf-8") as f:
                    rtf = f.read()
                raw_text = self.processor.clean_text(striprtf.rtf_to_text(rtf))
            else:
                raise ValueError(f"Unsupported file type: {ext}")
            
            # Calculate metrics and validate
            metrics = self.processor.calculate_metrics(raw_text)
            is_valid, issues = self.processor.validate_extraction_quality(raw_text, path)
            
            return {
                'text': raw_text,
                'metrics': metrics,
                'is_valid': is_valid,
                'issues': issues,
                'file_type': ext,
                'file_path': path
            }
            
        except Exception as e:
            logger.error(f"Error processing {path}: {e}")
            return {
                'text': '',
                'metrics': TextMetrics(0, 0, 0, 0.0, 0.0, 0.0, 0.0),
                'is_valid': False,
                'issues': [f"Extraction failed: {str(e)}"],
                'file_type': ext,
                'file_path': path
            }


# Backward compatibility functions
def ingest_to_text(path: str) -> str:
    """Backward compatible function - returns only text"""
    extractor = EnhancedCVExtractor()
    result = extractor.ingest_to_text_enhanced(path)
    return result['text']


if __name__ == '__main__':
    import argparse
    import json

    parser = argparse.ArgumentParser(description='Enhanced CV text extractor with semantic consistency.')
    parser.add_argument('input', help='Path or URL of the document')
    parser.add_argument('-o', '--output', help='Output text file (optional)')
    parser.add_argument('--metrics', action='store_true', help='Show extraction metrics')
    parser.add_argument('--validate', action='store_true', help='Validate extraction quality')
    
    args = parser.parse_args()
    
    extractor = EnhancedCVExtractor()
    result = extractor.ingest_to_text_enhanced(args.input)
    
    if args.metrics or args.validate:
        print(f"File: {result['file_path']}")
        print(f"File Type: {result['file_type']}")
        print(f"Valid: {'✅' if result['is_valid'] else '❌'}")
        
        if result['issues']:
            print("Issues found:")
            for issue in result['issues']:
                print(f"  - {issue}")
        
        if args.metrics:
            metrics = result['metrics']
            print(f"\nMetrics:")
            print(f"  Characters: {metrics.char_count:,}")
            print(f"  Words: {metrics.word_count:,}")
            print(f"  Lines: {metrics.line_count:,}")
            print(f"  Avg Line Length: {metrics.avg_line_length:.1f}")
            print(f"  Whitespace Ratio: {metrics.whitespace_ratio:.2%}")
            print(f"  Special Char Ratio: {metrics.special_char_ratio:.2%}")
        
        print("\n" + "="*50 + "\n")
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as fout:
            fout.write(result['text'])
        print(f"Text written to {args.output}")
    else:
        print(result['text'])